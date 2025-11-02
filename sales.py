import streamlit as st
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from scipy.interpolate import PchipInterpolator
import io
from google.oauth2 import service_account
from google.cloud import firestore
import base64
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from weasyprint import HTML
import plotly.express as px
import google.generativeai as genai
import vertexai
from vertexai.generative_models import GenerativeModel, Tool
from google.oauth2 import service_account
# --- Page Config ---
st.set_page_config(layout="wide", page_title="Advanced Business Plan Dashboard")
sns.set_theme(style="darkgrid", font_scale=1.1, palette="viridis")
# --- Global Settings ---
MODEL_START_YEAR = 2025

# <<< החלף את פונקציית ה-Stacked Bar שלך בזו המעודכנת >>>

# פונקציה 1: תרשים עמודות מוערם (Matplotlib/Seaborn) עם צבעים מותאמים אישית
def create_stacked_bar_chart(df):
    # !!! שינוי כאן: חזרה לעיצוב הכללי והסרת צבעים ספציפיים !!!
    sns.set_theme(style="darkgrid", font_scale=1.1, palette="viridis")
    fig, ax = plt.subplots(figsize=(16, 9))
    
    # --- !!! שינוי כאן: הסרת מילון הצבעים והמיון !!! ---
    # פשוט מציירים את כל הדאטה פריים
    df.plot(
        kind='bar', 
        stacked=True, 
        ax=ax, 
        # color=...  <-- הסרנו את הצבעים הקבועים
        width=0.7
    )
    
    # שאר הקוד של הפונקציה נשאר זהה
    for container in ax.containers:
        labels = [f'${v/1_000_000:.1f}M' if v > sum(df.sum())*0.015 else '' for v in container.datavalues]
        ax.bar_label(container, labels=labels, label_type='center', color='white', weight='bold', fontsize=10)

    totals = df.sum(axis=1)
    for i, total in enumerate(totals):
        if total > 0:
            ax.text(i, total + (totals.max() * 0.01), f'${total:,.0f}', ha='center', va='bottom', weight='bold', fontsize=12)

    ax.set_title('Total Sales Breakdown by Product', fontsize=20, weight='bold', pad=20)
    ax.set_ylabel('Revenue ($)', fontsize=14)
    ax.set_xlabel('Year', fontsize=14)
    ax.get_yaxis().set_major_formatter(plt.FuncFormatter(lambda x, p: f"${x/1_000_000:.0f}M"))
    ax.tick_params(axis='x', rotation=0, labelsize=12)
    ax.tick_params(axis='y', labelsize=12)
    ax.legend(title='Product', fontsize=12)
    ax.spines[['top', 'right']].set_visible(False)
    
    return fig

# פונקציה 2: תרשים עמודות מקובץ (Matplotlib/Seaborn)
def create_grouped_bar_chart(df):
    sns.set_theme(style="white")
    fig, ax = plt.subplots(figsize=(16, 9))
    
    df_melted = df.reset_index().rename(columns={'index': 'Year'}).melt(id_vars='Year', var_name='Product', value_name='Revenue')
    
    barplot = sns.barplot(data=df_melted, x='Year', y='Revenue', hue='Product', ax=ax, palette="viridis")
    
    for container in ax.containers:
        ax.bar_label(container, fmt='$ {:,.0f}', rotation=45, padding=5, fontsize=9, color='black')

    ax.set_title('Total Sales by Product (Grouped)', fontsize=20, weight='bold', pad=20)
    ax.set_ylabel('Revenue ($)', fontsize=14)
    ax.set_xlabel('Year', fontsize=14)
    ax.get_yaxis().set_major_formatter(plt.FuncFormatter(lambda x, p: f"${x/1_000_000:.0f}M"))
    ax.tick_params(axis='x', rotation=0, labelsize=12)
    ax.tick_params(axis='y', labelsize=12)
    ax.legend(title='Product', fontsize=12)
    ax.spines[['top', 'right']].set_visible(False)
    
    return fig

# פונקציה 3: תרשים אינטראקטיבי (Plotly)
def create_interactive_plotly_chart(df):
    df_melted = df.reset_index().rename(columns={'index': 'Year'}).melt(id_vars='Year', var_name='Product', value_name='Revenue')
    fig = px.bar(
        df_melted, 
        x='Year', 
        y='Revenue', 
        color='Product',
        barmode='group', # או 'stack' אם תעדיף
        title="Interactive Sales Breakdown by Product",
        labels={'Revenue': 'Revenue ($)', 'Year': 'Year'},
        template='plotly_white', # עיצוב נקי
        color_discrete_sequence=px.colors.sequential.Plasma_r, # פלטת צבעים יפה
        text_auto='.2s' # פורמט אוטומטי למספרים (למשל 4.5M)
    )
    fig.update_layout(
        title={'font': {'size': 22}, 'x': 0.5, 'xanchor': 'center'},
        legend_title_text='Product',
        height=600
    )
    fig.update_traces(textangle=0, textposition='outside')
    return fig
# פונקציית עזר להמרת גרף לתמונה שניתן להטמיע ב-HTML
def fig_to_base64_uri(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', bbox_inches='tight')
    buf.seek(0)
    base64_img = base64.b64encode(buf.read()).decode('utf-8')
    plt.close(fig) # סגירת הגרף לשחרור זיכרון
    return f'data:image/png;base64,{base64_img}'

# הפונקציה הראשית ליצירת ה-PDF
def to_pdf(results_dict):
    # --- 1. הגדרת עיצוב (CSS) - חזרה לגרסה הפשוטה ---
    html_style = """
    <style>
        @page { size: A4 portrait; margin: 1cm; }
        body { font-family: DejaVu Sans, Arial, sans-serif; }
        h1, h2, h3 { color: #003366; border-bottom: 2px solid #003366; padding-bottom: 5px;}
        h1 { font-size: 24pt; text-align: center; margin-bottom: 40px;}
        h2 { font-size: 18pt; margin-top: 50px;}
        h3 { font-size: 14pt; color: #335577; border-bottom: 1px solid #cccccc;}
        table { border-collapse: collapse; width: 100%; margin-top: 15px; margin-bottom: 25px; table-layout: auto; }
        th, td { border: 1px solid #dddddd; text-align: center; padding: 4px; font-size: 7pt;}
        th { background-color: #f2f2f2; font-weight: bold; }
        img { max-width: 100%; height: auto; display: block; margin-left: auto; margin-right: auto; margin-top: 15px; margin-bottom: 25px; }
        .page-break { page-break-before: always; }
    </style>
    """

    # --- 2. בניית גוף ה-HTML (גרסה מלאה ושלמה) ---
    html_body = f"<h1>Business Plan Analysis Report</h1><p style='text-align:center;'>Generated on: {pd.Timestamp.now(tz='Asia/Jerusalem').strftime('%d/%m/%Y')}</p>"
    
    product_list = [p for p in results_dict.keys() if p != 'summary']
    float_formatter = '{:,.0f}'.format
    main_display_start_date = pd.Timestamp('2026-01-01')

    # --- תוכן מלא עבור כל מוצר ---
    for product_name in product_list:
        data = results_dict[product_name]
        html_body += f"<div class='page-break'></div><h2>Analysis for: {product_name}</h2>"
        
        # טבלה 0 + גרף 0
        df_leads_q = data['lead_plan'][data['lead_plan'].index >= pd.Timestamp('2025-01-01')].T
        df_leads_q.columns = [f"{c.year}-Q{c.quarter}" for c in df_leads_q.columns]
        html_body += f"<h3>Table 0: Recommended Lead Contact Plan</h3>{df_leads_q.to_html(classes='dataframe', float_format=float_formatter)}"
        fig0 = create_yearly_bar_chart(data["lead_plan"][data["lead_plan"].index.year != 2030], "", "")
        html_body += f"<img src='{fig_to_base64_uri(fig0)}'>"
        
        # טבלה 1 + גרף 1
        df_acquired_q = data['acquired_customers_plan'][data['acquired_customers_plan'].index >= main_display_start_date].T
        df_acquired_q.columns = [f"{c.year}-Q{c.quarter}" for c in df_acquired_q.columns]
        html_body += f"<h3>Table 1: Acquired New Customers</h3>{df_acquired_q.to_html(classes='dataframe', float_format=float_formatter)}"
        fig1 = create_yearly_bar_chart(data['acquired_customers_plan'], "", "")
        html_body += f"<img src='{fig_to_base64_uri(fig1)}'>"
        
        # טבלה 2 + גרף 2
        df_cum_q = data['cumulative_customers'][data['cumulative_customers'].index >= main_display_start_date].T
        df_cum_q.columns = [f"{c.year}-Q{c.quarter}" for c in df_cum_q.columns]
        html_body += f"<h3>Table 2: Cumulative Customers</h3>{df_cum_q.to_html(classes='dataframe', float_format=float_formatter)}"
        fig2 = create_yearly_bar_chart(data['cumulative_customers'], "", "", is_cumulative=True)
        html_body += f"<img src='{fig_to_base64_uri(fig2)}'>"

        # טבלה 3 + גרף 3
        html_body += f"<h3>Table 3: Target vs. Actual Revenue</h3>{data['validation_df'].to_html(classes='dataframe', float_format=float_formatter)}"
        plot_df_melted = data['validation_df'].reset_index().melt(id_vars='Year', var_name='Type', value_name='Revenue')
        fig3, ax3 = plt.subplots(figsize=(10, 5))
        sns.barplot(data=plot_df_melted, x='Year', y='Revenue', hue='Type', ax=ax3, palette="mako")
        for c in ax3.containers: ax3.bar_label(c, fmt='${:,.0f}', padding=3, fontsize=8)
        html_body += f"<img src='{fig_to_base64_uri(fig3)}'>"
        
        # טבלאות 4 ו-5
        html_body += "<h3>Underlying Assumptions</h3>"
        # !!! שינוי כאן: טונות -> יחידות !!!
        html_body += "<h4>Table 4: Annual Units per Single Customer</h4>" + data['units_per_customer'].to_html(classes='dataframe', float_format='{:.2f}'.format)
        html_body += "<h4>Table 5: Generated Penetration Rates (%)</h4>" + (data['pen_rate_df'] * 100).to_html(classes='dataframe', float_format='{:.1f}%'.format)

    # --- תוכן מלא עבור הסיכום הכללי ---
    summary_data = results_dict.get("summary", {})
    if summary_data:
        html_body += "<div class='page-break'></div><h2>Overall Summary</h2>"
        
        # טבלת סיכום הכנסות
        html_body += "<h3>Total Revenue per Year</h3>" + summary_data["summary_revenue"].to_html(classes='dataframe', float_format=float_formatter)
        
        # טבלת לקוחות מצטברת של הסיכום
        summary_customers_to_display = summary_data["summary_customers_raw"][summary_data["summary_customers_raw"].index >= main_display_start_date].to_frame("Total Customers").T
        summary_customers_to_display.columns = [f"{c.year}-Q{c.quarter}" for c in summary_customers_to_display.columns]
        html_body += f"<h3>Total Cumulative Customers (Quarterly)</h3>{summary_customers_to_display.to_html(classes='dataframe', float_format=float_formatter)}"

        # גרף סיכום הכנסות - החלק שהיה חסר
        all_revenues = {p: results_dict[p]['annual_revenue'] for p in product_list}
        summary_plot_df = pd.DataFrame(all_revenues)
        summary_plot_df_melted = summary_plot_df.reset_index().rename(columns={'index': 'Year'}).melt(id_vars='Year', var_name='Product', value_name='Revenue')
        fig_sum, ax_sum = plt.subplots(figsize=(10, 5))
        summary_barplot = sns.barplot(data=summary_plot_df_melted, x='Year', y='Revenue', hue='Product', ax=ax_sum, palette="rocket_r")
        for container in ax_sum.containers: ax_sum.bar_label(container, fmt='$ {:,.0f}', rotation=45, padding=8, fontsize=8, color='black', fontweight='bold')
        html_body += f"<h3>Total Revenue Breakdown by Product</h3><img src='{fig_to_base64_uri(fig_sum)}'>"

    # --- 3. הרכבת ה-HTML המלא ויצירת ה-PDF ---
    full_html = f"<!DOCTYPE html><html><head><meta charset='UTF-8'><title>Report</title>{html_style}</head><body>{html_body}</body></html>"    
    pdf_bytes = HTML(string=full_html).write_pdf()
    return pdf_bytes

def add_fig_to_slide(slide, fig, left, top, width):
    """
    Saves a matplotlib figure to a buffer and adds it correctly and safely to a slide.
    This helper function prevents the UnidentifiedImageError.
    """
    img_buffer = io.BytesIO()
    fig.savefig(img_buffer, format='png', bbox_inches='tight')
    img_buffer.seek(0)
    slide.shapes.add_picture(img_buffer, left, top, width=width)
    plt.close(fig) # Close the figure to free up memory

def add_df_to_slide(slide, df, left, top, width, height, font_size=9):
    """
    Helper function to add a pandas DataFrame to a PowerPoint slide with better formatting.
    """
    df_title = df.name if hasattr(df, 'name') else "Data Table"
    slide.shapes.add_textbox(left, top, width, Inches(0.4)).text_frame.text = df_title
    top += Inches(0.4)
    rows, cols = df.shape
    rows += 1
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    for i in range(cols):
        table.columns[i].width = Inches(width.inches / cols)
    for i, col_name in enumerate(df.columns):
        cell = table.cell(0, i)
        cell.text = str(col_name)
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(font_size)
        p.alignment = PP_ALIGN.CENTER
    for r in range(rows - 1):
        for c in range(cols):
            cell = table.cell(r + 1, c)
            value = df.iloc[r, c]
            if isinstance(value, (int, float)):
                cell.text = f"{value:,.0f}"
            else:
                cell.text = str(value)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(font_size - 1)
            p.alignment = PP_ALIGN.CENTER
    return table_shape

def create_product_presentation(product_name, data):
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # TODO: עדכן את המספרים האלה לפי הפלט של check_layouts.py
    title_slide_layout = prs.slide_layouts[0] 
    blank_slide_layout = prs.slide_layouts[6]

    # --- Title Slide ---
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = f"Business Plan Analysis: {product_name}"

    # =======================================================
    #                 *** START OF THE FIX ***
    # This will now safely handle layouts with no subtitle
    # =======================================================
    try:
        current_date = pd.Timestamp.now(tz='Asia/Jerusalem').strftime('%d/%m/%Y')
        slide.placeholders[1].text = f"Generated on: {current_date}"
    except (KeyError, IndexError):
        # This layout doesn't have a subtitle placeholder, which is fine.
        pass
    # =======================================================
    #                 *** END OF THE FIX ***
    # =======================================================
        
    # ... (שאר הקוד של הפונקציה ממשיך כרגיל) ...
    df_leads_q = data['lead_plan'].T
    df_leads_q.columns = [f"{c.year}-Q{c.quarter}" for c in df_leads_q.columns]
    df_acquired_q = data['acquired_customers_plan'].T
    df_acquired_q.columns = [f"{c.year}-Q{c.quarter}" for c in df_acquired_q.columns]
    df_cumulative_q = data['cumulative_customers'].T
    df_cumulative_q.columns = [f"{c.year}-Q{c.quarter}" for c in df_cumulative_q.columns]
    slide = prs.slides.add_slide(blank_slide_layout)
    fig = create_yearly_bar_chart(data['lead_plan'][data['lead_plan'].index.year != 2030], "Chart 0: Leads to Contact per Year", "")
    add_fig_to_slide(slide, fig, Inches(1), Inches(1), width=Inches(14))
    slide = prs.slides.add_slide(blank_slide_layout)
    df_leads_q.name = "Table 0: Recommended Lead Contact Plan (Quarterly)"
    add_df_to_slide(slide, df_leads_q, Inches(0.5), Inches(0.2), Inches(15), Inches(3))
    slide = prs.slides.add_slide(blank_slide_layout)
    fig = create_yearly_bar_chart(data['acquired_customers_plan'], "Chart 1: Acquired New Customers per Year", "")
    add_fig_to_slide(slide, fig, Inches(1), Inches(1), width=Inches(14))
    slide = prs.slides.add_slide(blank_slide_layout)
    df_acquired_q.name = "Table 1: Acquired New Customers (Quarterly)"
    add_df_to_slide(slide, df_acquired_q, Inches(0.5), Inches(0.2), Inches(15), Inches(3))
    slide = prs.slides.add_slide(blank_slide_layout)
    fig = create_yearly_bar_chart(data['cumulative_customers'], "Chart 2: Cumulative Customers at Year End", "", is_cumulative=True)
    add_fig_to_slide(slide, fig, Inches(1), Inches(1), width=Inches(14))
    slide = prs.slides.add_slide(blank_slide_layout)
    df_cumulative_q.name = "Table 2: Cumulative Customers (Quarterly)"
    add_df_to_slide(slide, df_cumulative_q, Inches(0.5), Inches(0.2), Inches(15), Inches(3))
    slide = prs.slides.add_slide(blank_slide_layout)
    slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(15), Inches(0.8)).text_frame.text = "Underlying Assumptions"
    
    # !!! שינוי כאן: טונות -> יחידות !!!
    df_units = data['units_per_customer'].T
    df_units.name = "Table 4: Annual Units per Single Customer"
    add_df_to_slide(slide, df_units.style.format("{:,.2f}").data, Inches(0.5), Inches(1), Inches(15), Inches(2.5), font_size=12)
    
    df_pen = (data['pen_rate_df'] * 100).T
    df_pen.name = "Table 5: Generated Penetration Rates (%)"
    add_df_to_slide(slide, df_pen.style.format("{:,.1f}%").data, Inches(0.5), Inches(4), Inches(15), Inches(2.5), font_size=12)
    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer.getvalue()
def create_summary_presentation(summary_data, all_results):
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # TODO: עדכן את המספרים האלה לפי התבנית שלך!
    title_slide_layout = prs.slide_layouts[0]
    blank_slide_layout = prs.slide_layouts[6]

    # --- Title Slide ---
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Overall Summary Report"
    
    # =======================================================
    #                 *** START OF THE FIX ***
    # This will now safely handle layouts with no subtitle
    # =======================================================
    try:
        current_date = pd.Timestamp.now(tz='Asia/Jerusalem').strftime('%d/%m/%Y')
        slide.placeholders[1].text = f"Generated on: {current_date}"
    except (KeyError, IndexError):
        # This layout doesn't have a subtitle placeholder, which is fine.
        pass
    # =======================================================
    #                 *** END OF THE FIX ***
    # =======================================================
        
    # ... (שאר הקוד של הפונקציה ממשיך כרגיל) ...
    slide = prs.slides.add_slide(blank_slide_layout)
    slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(15), Inches(0.8)).text_frame.text = "Total Revenue Breakdown by Product"
    product_list = [p for p in all_results.keys() if p != 'summary']
    all_revenues = {p: all_results[p]['annual_revenue'] for p in product_list}
    summary_plot_df = pd.DataFrame(all_revenues)
    summary_plot_df_melted = summary_plot_df.reset_index().rename(columns={'index': 'Year'}).melt(id_vars='Year', var_name='Product', value_name='Revenue')
    fig, ax = plt.subplots(figsize=(12, 6))
    barplot = sns.barplot(data=summary_plot_df_melted, x='Year', y='Revenue', hue='Product', ax=ax, palette="rocket_r")
    ax.get_yaxis().set_major_formatter(plt.FuncFormatter(lambda x, p: f"${x/1_000_000:.0f}M"))
    for container in barplot.containers:
        ax.bar_label(container, fmt='$ {:,.0f}', rotation=45, padding=8, fontsize=10, color='black', fontweight='bold')
    add_fig_to_slide(slide, fig, Inches(1), Inches(1.2), width=Inches(14))
    slide = prs.slides.add_slide(blank_slide_layout)
    df_summary_cust = summary_data["summary_customers_raw"].to_frame("Total Customers").T
    df_summary_cust.columns = [f"{c.year}-Q{c.quarter}" for c in df_summary_cust.columns]
    df_summary_cust.name = "Total Cumulative Customers (Quarterly)"
    add_df_to_slide(slide, df_summary_cust, Inches(0.5), Inches(1.5), Inches(15), Inches(2))
    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer.getvalue()
# =========================
# פונקציות המרה ל/מ Firestore
# =========================
def serialize_for_firestore(value):
    """המרה לאובייקט שניתן לשמור ב־Firestore וב־session_state"""
    if isinstance(value, pd.DataFrame):
        return {"__type__": "DataFrame", "data": value.to_dict(orient='split')}
    elif isinstance(value, pd.Series):
        return {"__type__": "Series", "data": value.to_dict()}
    elif isinstance(value, pd.Timestamp):
        return {"__type__": "Timestamp", "data": value.isoformat()}
    elif isinstance(value, bytes):
        return {"__type__": "Bytes", "data": base64.b64encode(value).decode('utf-8')}
    elif isinstance(value, (str, int, float, bool, type(None))):
        return value
    elif isinstance(value, list):
        return [serialize_for_firestore(v) for v in value]
    elif isinstance(value, dict):
        return {k: serialize_for_firestore(v) for k, v in value.items()}
    else:
        return {"__type__": "str", "data": str(value)}

def deserialize_from_firestore(value):
    """שחזור הערך לסוג המקורי לאחר טעינה"""
    if isinstance(value, dict) and "__type__" in value:
        t = value["__type__"]
        if t == "DataFrame":
            return pd.DataFrame(**value["data"])
        elif t == "Series":
            return pd.Series(value["data"])
        elif t == "Timestamp":
            return pd.Timestamp(value["data"])
        elif t == "Bytes":
            return base64.b64decode(value["data"])
        elif t == "str":
            return value["data"]
    if isinstance(value, list):
        return [deserialize_from_firestore(v) for v in value]
    if isinstance(value, dict):
        return {k: deserialize_from_firestore(v) for k, v in value.items()}
    return value

# --- Session State Initialization ---
if 'products' not in st.session_state:
    # !!! שינוי כאן: מוצרים גנריים כברירת מחדל !!!
    st.session_state.products = ["Product 1"]
if 'results' not in st.session_state:
    st.session_state.results = {}

# --- Excel Export ---
@st.cache_data
# <<< החלף את כל פונקציית to_excel שלך בקוד הבא >>>

def to_excel(results_dict):
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine='openpyxl') as writer:
        # Loop for each product to create a dedicated sheet
        for product_name, data in results_dict.items():
            
            # --- !!! התיקון נמצא כאן !!! ---
            # מדלגים על ריצה אם שם המוצר ריק או אם זה מפתח הסיכום
            if not product_name or product_name == 'summary':
                continue
            # --- סוף התיקון ---
            
            # --- Prepare all dataframes for the sheet ---
            # שימוש ב-.get() כדי למנוע קריסות אם מפתח חסר
            df_lead_plan_T = data.get('lead_plan', pd.DataFrame()).T
            df_acquired_cust_T = data.get('acquired_customers_plan', pd.DataFrame()).T
            df_cum_cust_q_T = data.get("cumulative_customers", pd.DataFrame()).T
            df_validation = data.get('validation_df', pd.DataFrame())
            
            # !!! שינוי כאן: טונות -> יחידות !!!
            df_units_per_customer = data.get('units_per_customer', pd.DataFrame()).T
            df_pen_rate = (data.get('pen_rate_df', pd.DataFrame()) * 100).T
            
            # הוספת הטבלאות החדשות
            df_revenue_by_type_T = data.get('revenue_by_type_q', pd.DataFrame()).T
            # !!! שינוי כאן: טונות -> יחידות !!!
            df_units_by_type_T = data.get('units_by_type_q', pd.DataFrame()).T


            # Format all quarterly columns
            quarterly_dfs = [
                df_lead_plan_T, df_acquired_cust_T, df_cum_cust_q_T, 
                df_revenue_by_type_T, df_units_by_type_T # <-- שינוי כאן
            ]
            for df in quarterly_dfs:
                if not df.empty:
                    df.columns = [f"{c.year}-Q{c.quarter}" for c in df.columns]

            # --- Write tables to the sheet one by one with correct ordering ---
            current_row = 1
            
            # Table 0: Leads
            df_lead_plan_T.to_excel(writer, sheet_name=product_name, startrow=current_row + 1)
            writer.sheets[product_name].cell(row=current_row, column=1, value="Recommended Lead Contact Plan (Table 0)")
            current_row += df_lead_plan_T.shape[0] + 5

            # Table A: Revenue by Type (NEW)
            df_revenue_by_type_T.to_excel(writer, sheet_name=product_name, startrow=current_row + 1)
            writer.sheets[product_name].cell(row=current_row, column=1, value="Revenue per Customer Type ($) (Table A)")
            current_row += df_revenue_by_type_T.shape[0] + 5

            # !!! שינוי כאן: טונות -> יחידות !!!
            # Table B: Units by Type (NEW)
            df_units_by_type_T.to_excel(writer, sheet_name=product_name, startrow=current_row + 1)
            writer.sheets[product_name].cell(row=current_row, column=1, value="Units Sold per Customer Type (Table B)")
            current_row += df_units_by_type_T.shape[0] + 5

            # Table 1: Acquired Customers
            df_acquired_cust_T.to_excel(writer, sheet_name=product_name, startrow=current_row + 1)
            writer.sheets[product_name].cell(row=current_row, column=1, value="Acquired New Customers per Quarter (Table 1)")
            current_row += df_acquired_cust_T.shape[0] + 5

            # Table 2: Cumulative Customers
            df_cum_cust_q_T.to_excel(writer, sheet_name=product_name, startrow=current_row + 1)
            writer.sheets[product_name].cell(row=current_row, column=1, value="Cumulative Customers (Quarterly) (Table 2)")
            current_row += df_cum_cust_q_T.shape[0] + 5
            
            # Table 3: Validation
            df_validation.to_excel(writer, sheet_name=product_name, startrow=current_row + 1)
            writer.sheets[product_name].cell(row=current_row, column=1, value="Target vs. Actual Revenue (Table 3)")
            current_row += df_validation.shape[0] + 5

            # !!! שינוי כאן: טונות -> יחידות !!!
            # Table 4: Units per Customer
            df_units_per_customer.to_excel(writer, sheet_name=product_name, startrow=current_row + 1)
            writer.sheets[product_name].cell(row=current_row, column=1, value="Annual Units per Single Customer (Target-Driven) (Table 4)")
            current_row += df_units_per_customer.shape[0] + 5

            # Table 5: Penetration Rate
            df_pen_rate.to_excel(writer, sheet_name=product_name, startrow=current_row + 1)
            writer.sheets[product_name].cell(row=current_row, column=1, value="Generated Penetration Rates to Meet Target (%) (Table 5)")

        # --- Overall Summary Sheet ---
        if "summary" in results_dict:
            summary_data = results_dict["summary"]
            if summary_data and summary_data.get("summary_revenue") is not None:
                summary_revenue_df = summary_data["summary_revenue"]
                summary_customers_df = summary_data["summary_customers_raw"]
                summary_revenue_df.to_excel(writer, sheet_name="Overall Summary", startrow=2)
                writer.sheets["Overall Summary"].cell(row=1, column=1, value="Total Revenue per Year")
                
                summary_customers_df_T = summary_customers_df.to_frame("Total Customers").T
                if not summary_customers_df_T.empty:
                    summary_customers_df_T.columns = [f"{c.year}-Q{c.quarter}" for c in summary_customers_df_T.columns]
                
                summary_customers_df_T.to_excel(writer, sheet_name="Overall Summary", startrow=10)
                writer.sheets["Overall Summary"].cell(row=9, column=1, value="Total Cumulative Customers (Quarterly)")
            
    return output.getvalue()
    # --- END OF THE FINAL CORRECTED FUNCTION ---
# --- Firebase ---
@st.cache_resource
def init_connection():
    try:
        creds_json = dict(st.secrets.firebase)
        creds = service_account.Credentials.from_service_account_info(creds_json)
        return firestore.Client(credentials=creds, project=creds_json['project_id'])
    except Exception as e:
        st.error(f"Failed to connect to Firebase. Error: {e}")
        return None

db = init_connection()

# --- Save/Load (עכשיו עם סריאליזציה) ---
def save_scenario(user_id, scenario_name, data):
    if not db or not user_id or not scenario_name:
        st.sidebar.warning("User ID and Scenario Name are required to save.")
        return
    try:
        data_to_save = {}
        for k, v in data.items():
            # --- START OF FIX ---
            # Added 'pie_select' to the list of prefixes to ignore during save
            if isinstance(k, str) and not k.startswith(('results', 'load_scenario_select', 'scenario_name', 'new_product_name_input', 'confirm_delete_checkbox', 'FormSubmitter', '_', 'pie_select')):
            # --- END OF FIX ---
                data_to_save[k] = serialize_for_firestore(v)
        
        db.collection('users').document(user_id).collection('scenarios').document(scenario_name).set(data_to_save)
        st.sidebar.success(f"Scenario '{scenario_name}' saved!")
    except Exception as e:
        st.sidebar.error(f"Error saving scenario: {e}")
def create_yearly_bar_chart(df_quarterly, title, y_axis_label, is_cumulative=False):
    """
    Creates a yearly grouped bar chart from quarterly data.
    - If is_cumulative is True, it takes the last value of each year.
    - Otherwise, it sums the values for each year.
    """
    # 1. Data Preparation
    if is_cumulative:
        # For cumulative data, we want the value at the end of each year
        df_yearly = df_quarterly.resample('YE').last()
    else:
        # For new leads/customers, we sum the quarters to get the yearly total
        df_yearly = df_quarterly.resample('YE').sum()
    
    df_yearly.index = df_yearly.index.year # Use just the year number for the x-axis
    df_yearly.index.name = "Year"
    
    # Melt the dataframe to make it "tidy" for Seaborn
    df_melted = df_yearly.reset_index().melt(
        id_vars='Year', 
        var_name='Customer Type', 
        value_name='Count'
    )
    
    # 2. Plotting
    fig, ax = plt.subplots(figsize=(14, 7))
    sns.barplot(data=df_melted, x='Year', y='Count', hue='Customer Type', ax=ax, palette='viridis')
    
    # 3. Aesthetics
    ax.set_title(title, fontsize=18, weight='bold', pad=20)
    ax.set_xlabel("Year", fontsize=12)
    ax.set_ylabel(y_axis_label, fontsize=12)
    ax.get_yaxis().set_major_formatter(plt.FuncFormatter(lambda x, p: format(int(x), ',')))
    ax.legend(title='Customer Type')
    
    # Add labels on top of each bar
    for container in ax.containers:
        ax.bar_label(container, fmt='{:,.0f}', padding=3, fontsize=9)
        
    plt.tight_layout()
    return fig
    
def get_user_scenarios(user_id):
    if not db or not user_id:
        return []
    try:
        docs = db.collection('users').document(user_id).collection('scenarios').stream()
        return [""] + [doc.id for doc in docs]
    except Exception as e:
        st.sidebar.error(f"Error fetching scenarios: {e}")
        return [""]
def delete_scenario(user_id, scenario_name):
    """Deletes a specific scenario for a user from Firestore."""
    if not db or not user_id or not scenario_name:
        st.sidebar.warning("Could not delete scenario. User ID or scenario name is missing.")
        return False
    try:
        db.collection('users').document(user_id).collection('scenarios').document(scenario_name).delete()
        st.sidebar.success(f"Scenario '{scenario_name}' deleted successfully.")
        return True
    except Exception as e:
        st.sidebar.error(f"Error deleting scenario: {e}")
        return False

def load_scenario_data(user_id, scenario_name):
    if not db or not user_id or not scenario_name:
        return None
    try:
        doc_ref = db.collection('users').document(user_id).collection('scenarios').document(scenario_name)
        doc = doc_ref.get()
        if doc.exists:
            st.sidebar.info(f"Loaded '{scenario_name}'.")
            return doc.to_dict()
        else:
            st.sidebar.warning("Scenario not found.")
            return None
    except Exception as e:
        st.sidebar.error(f"Error loading: {e}")
        return None

# <<< פונקציית calculate_plan מתוקנת ומלאה >>>

# <<< החלף את כל פונקציית calculate_plan שלך בקוד הבא >>>

# <<< החלף את כל פונקציית calculate_plan שלך בקוד הבא >>>

def calculate_plan(is_m, is_l, is_g, market_gr, pen_y1, tt_m, tt_l, tt_g, 
                   # !!! שינוי כאן: טונות -> יחידות !!!
                   annual_rev_targets, f_m, f_l, f_g, ip_unit, pdr, price_floor,
                   cost_quantities_u, cost_values_per_unit,
                   global_start_year, global_start_quarter, launch_year): # <<< עדכון חתימת הפונקציה

    # הגדרות כלליות
    MODEL_START_YEAR = 2025
    CALCULATION_START_YEAR = MODEL_START_YEAR
    NUM_YEARS = 6
    years = np.array([CALCULATION_START_YEAR + i for i in range(NUM_YEARS)])
    quarters_index = pd.date_range(start=f'{CALCULATION_START_YEAR}-01-01', periods=NUM_YEARS*4, freq='QE')
    customer_types = ['Medium', 'Large', 'Global']
    
    # --- חלק 1: חישוב מנועי הערך (עם לוגיקת השקה דינמית) ---
    
    # !!! שינוי כאן: טונות -> יחידות !!!
    units_per_customer = pd.DataFrame(0.0, index=years, columns=customer_types, dtype=float)
    
    last_model_year = years[-1]
    
    # אם שנת ההשקה היא אחרי תקופת המודל, כל החישובים יהיו אפס
    if launch_year <= last_model_year:
        # הצבת הערך ההתחלתי בשנה הנכונה
        units_per_customer.loc[launch_year] = [is_m, is_l, is_g]

        # חישוב תקופת המכירה האמיתית של המוצר
        sales_duration_years = last_model_year - launch_year + 1
        
        # חישוב אחוזי חדירה דינמיים
        # !!! שינוי כאן: טונות -> יחידות !!!
        initial_units = {'Medium': is_m, 'Large': is_l, 'Global': is_g}
        target_units = {'Medium': tt_m, 'Large': tt_l, 'Global': tt_g}
        
        # אינדקס יחסי לתקופת המכירה
        relative_year_index = range(1, sales_duration_years + 1)
        pen_rate_df_relative = pd.DataFrame(index=relative_year_index, columns=customer_types)

        for c_type in customer_types:
            total_market_growth_factor = (1 + market_gr / 100) ** (sales_duration_years - 1)
            # !!! שינוי כאן: טונות -> יחידות !!!
            if initial_units[c_type] == 0: required_pen_growth_factor = 1.0
            else: required_pen_growth_factor = (target_units[c_type] / initial_units[c_type]) / total_market_growth_factor
            
            pen_rate_y_final = (pen_y1 / 100) * required_pen_growth_factor
            
            # התאמת נקודות האינטרפולציה לתקופה הדינמית
            interp_points_x = [1, sales_duration_years / 2, sales_duration_years] if sales_duration_years > 2 else [1, sales_duration_years]
            interp_points_y = [pen_y1 / 100, (pen_y1/100 + pen_rate_y_final)/2, pen_rate_y_final] if sales_duration_years > 2 else [pen_y1 / 100, pen_rate_y_final]
            
            interp_func = PchipInterpolator(interp_points_x, interp_points_y)
            pen_rate_df_relative[c_type] = interp_func(relative_year_index)

        # הרצת לולאת הצמיחה רק על השנים הרלוונטיות
        launch_year_index_in_years_array = list(years).index(launch_year)
        for i in range(launch_year_index_in_years_array, NUM_YEARS - 1):
            current_year, prev_year = years[i+1], years[i]
            relative_idx_current, relative_idx_prev = i - launch_year_index_in_years_array + 2, i - launch_year_index_in_years_array + 1
            
            for c_type in customer_types:
                # !!! שינוי כאן: טונות -> יחידות !!!
                prev_units = units_per_customer.loc[prev_year, c_type]
                market_growth_factor = (1 + market_gr / 100)
                pen_growth_factor = pen_rate_df_relative.loc[relative_idx_current, c_type] / pen_rate_df_relative.loc[relative_idx_prev, c_type]
                units_per_customer.loc[current_year, c_type] = prev_units * market_growth_factor * pen_growth_factor
    
    pen_rate_df = pen_rate_df_relative if launch_year <= last_model_year else pd.DataFrame(0.0, index=range(1, NUM_YEARS + 1), columns=customer_types)

    # --- חישוב מחירים (ללא שינוי) ---
    prices = []
    # !!! שינוי כאן: ip_kg -> ip_unit !!!
    current_price = ip_unit
    decay_rate = pdr / 100.0
    for _ in quarters_index:
        prices.append(current_price)
        next_price = current_price * (1 - decay_rate)
        current_price = max(next_price, price_floor)
        
    # !!! שינוי כאן: הסרת * 1000. המחיר הוא עכשיו ליחידה !!!
    price_per_unit_q = pd.Series(prices, index=quarters_index)
    # !!! שינוי כאן: טונות -> יחידות !!!
    units_per_cust_q = units_per_customer.loc[quarters_index.year].set_axis(quarters_index) / 4

    # --- חלק 2: מנוע החישוב ההפוך (Top-Down) ---
    Q_GROWTH_RATE = 0.10
    growth_factors = np.array([1, (1 + Q_GROWTH_RATE), (1 + Q_GROWTH_RATE)**2, (1 + Q_GROWTH_RATE)**3])
    quarterly_weights = growth_factors / growth_factors.sum()
    
    quarterly_rev_targets_list = []
    for yearly_target in annual_rev_targets:
        quarterly_rev_targets_list.extend(yearly_target * quarterly_weights)
    
    quarterly_rev_targets = pd.Series(quarterly_rev_targets_list, index=quarters_index)
    
    # --- !!! שינוי מרכזי כאן: הגדרת תאריך התחלה אפקטיבי !!! ---
    global_start_date = pd.Timestamp(f"{global_start_year}-{(global_start_quarter-1)*3 + 1}-01")
    product_launch_date = pd.Timestamp(f"{launch_year}-01-01")
    
    effective_start_date = max(global_start_date, product_launch_date)
    
    quarterly_rev_targets.loc[quarterly_rev_targets.index < effective_start_date] = 0
    # --- סוף השינוי ---
    
    total_focus = f_m + f_l + f_g
    if total_focus == 0: return {"error": "Total Sales Focus must be greater than 0."}
    focus_norm = {'Medium': f_m / total_focus, 'Large': f_l / total_focus, 'Global': f_g / total_focus}
    
    new_customers_plan = pd.DataFrame(0.0, index=quarters_index, columns=customer_types)
    cumulative_customers = pd.DataFrame(0.0, index=quarters_index, columns=customer_types)
    
    for i, q_date in enumerate(quarters_index):
        # הלולאה ממשיכה כרגיל, היא תושפע מהיעדים המאופסים
        prev_cumulative = cumulative_customers.iloc[i-1] if i > 0 else pd.Series(0.0, index=customer_types)
        # !!! שינוי כאן: טונות -> יחידות !!!
        value_per_customer_type = units_per_cust_q.loc[q_date] * price_per_unit_q.loc[q_date]
        revenue_from_existing = (value_per_customer_type * prev_cumulative).sum()
        revenue_gap = quarterly_rev_targets.loc[q_date] - revenue_from_existing
        
        if revenue_gap > 0:
            blended_revenue_per_customer = (value_per_customer_type * pd.Series(focus_norm)).sum()
            if blended_revenue_per_customer > 0:
                total_new_customers_needed = revenue_gap / blended_revenue_per_customer
                for c_type in customer_types:
                    new_customers_plan.loc[q_date, c_type] = total_new_customers_needed * focus_norm[c_type]
                    
        cumulative_customers.loc[q_date] = prev_cumulative + new_customers_plan.loc[q_date]

    # --- חלק 3: חישוב פלטים סופיים ---
    # !!! שינוי כאן: טונות -> יחידות !!!
    revenue_per_customer_type_q = units_per_cust_q.mul(price_per_unit_q, axis=0)
    actual_revenue_q = (revenue_per_customer_type_q * cumulative_customers.round().astype(int)).sum(axis=1)

    # !!! שינוי כאן: טונות -> יחידות !!!
    units_by_type_q = units_per_cust_q * cumulative_customers.round().astype(int)
    # !!! שינוי כאן: תיקון לוגיקה + שם משתנה. המחיר הוא כבר ליחידה !!!
    revenue_by_type_q = units_by_type_q.mul(price_per_unit_q, axis=0)
    
    # !!! שינוי כאן: טונות -> יחידות !!!
    total_units_q = units_by_type_q.sum(axis=1)
    
    # !!! שינוי כאן: טונות -> יחידות !!!
    cost_per_unit_q = pd.Series(0.0, index=quarters_index)
    sorted_quantities = sorted(cost_quantities_u)
    sorted_values = [cost_values_per_unit[cost_quantities_u.index(q)] for q in sorted_quantities]

    for i, q_date in enumerate(quarters_index):
        # !!! שינוי כאן: טונות -> יחידות !!!
        units_sold = total_units_q.loc[q_date]
        if units_sold > 0:
            cost_idx = np.searchsorted(sorted_quantities, units_sold, side='right') -1
            # !!! שינוי כאן: cost_per_kg -> cost_per_unit !!!
            cost_per_unit = sorted_values[cost_idx] if cost_idx >= 0 else sorted_values[0]
            # !!! שינוי כאן: הסרת * 1000, העלות היא כבר ליחידה !!!
            cost_per_unit_q.loc[q_date] = cost_per_unit
    
    # !!! שינוי כאן: טונות -> יחידות !!!
    total_cost_q = total_units_q * cost_per_unit_q
    profit_q = actual_revenue_q - total_cost_q
    
    annual_revenue_series = actual_revenue_q.resample('YE').sum()
    
    return {
        "cumulative_customers": cumulative_customers,
        "annual_revenue": annual_revenue_series,
        "annual_revenue_targets": pd.Series(annual_rev_targets, index=years),
        # !!! שינוי כאן: טונות -> יחידות !!!
        "units_by_type_q": units_by_type_q,
        "revenue_by_type_q": revenue_by_type_q,
        "profit_q": profit_q,
        "total_production_cost_q": total_cost_q,
        # !!! שינוי כאן: טונות -> יחידות !!!
        "units_per_customer": units_per_customer,
        "pen_rate_df": pen_rate_df,
        "error": None
    }
    
    # --- הקוד הכפול והשבור שהיה כאן - נמחק ---
    
# ==================================
# --- פונקציות עזר עבור AI TOOLS ---
# ==================================

def update_setting(setting_key: str, new_value: any):
    """
    פונקציית פייתון שמעדכנת ערך ב-Streamlit's session_state.
    ה-AI יקרא לפונקציה הזו כדי לשנות הגדרות באפליקציה.
    """
    try:
        # ננסה להמיר את הערך לסוג הנכון (למשל, אם ה-AI שולח "25" במקום 25.0)
        original_value = st.session_state.get(setting_key)
        if original_value is not None:
            value_type = type(original_value)
            try:
                converted_value = value_type(new_value)
                st.session_state[setting_key] = converted_value
                return f"Success: Set {setting_key} to {converted_value}"
            except Exception as e:
                # אם ההמרה נכשלה, פשוט נכניס את הערך כמו שהוא
                st.session_state[setting_key] = new_value
                return f"Success (with fallback): Set {setting_key} to {new_value}. Conversion error: {e}"
        else:
            # אם המפתח לא קיים, צור אותו
            st.session_state[setting_key] = new_value
            return f"Success: Created and set {setting_key} to {new_value}"
            
    except Exception as e:
        return f"Error: Could not set {setting_key}. Reason: {e}"

# --- הגדרת הכלים (Tools) עבור Gemini ---
# זה ה"תפריט" שאומר ל-AI אילו פונקציות הוא יכול להפעיל
tools_schema = [
    {
        "name": "update_setting",
        "description": "עדכון הגדרה, פרמטר, או ערך קלט באפליקציה. השתמש בזה כדי לשנות מחירים, עלויות, כמויות, או כל דבר שהמשתמש מבקש לשנות.",
        "parameters": {
            "type": "OBJECT",
            "properties": {
                "setting_key": {
                    "type": "STRING",
                    "description": "שם המפתח המדויק ב-session_state, למשל 'ip_unit_Product 1' או 'start_year'."
                },
                "new_value": {
                    "type": "STRING", # שימוש ב-STRING מפשט את העבודה, נמיר סוגים בפייתון
                    "description": "הערך החדש שיש להגדיר עבור המפתח."
                }
            },
            "required": ["setting_key", "new_value"]
        }
    }
]

# מיפוי שם הפונקציה שה-AI קורא לה לפונקציית הפייתון האמיתית
available_tools = {
    "update_setting": update_setting,
}
# ==================================
# --- סוף פונקציות ה-AI ---
# ==================================
def create_lead_plan(acquired_customers_plan, success_rates, time_aheads_in_quarters):
    # --- START OF CHANGE: Lead plan now creates its own extended timeline ---
    LEAD_START_YEAR = 2025
    
    # אם אין לקוחות שצריך לגייס, החזר טבלה ריקה
    if acquired_customers_plan.empty:
        return pd.DataFrame()

    # יצירת ציר זמן מורחב עבור תוכנית הלידים, החל מ-2025 ועד סוף התוכנית הראשית
    lead_quarters_index = pd.date_range(
        start=f'{LEAD_START_YEAR}-01-01', 
        end=acquired_customers_plan.index.max(), 
        freq='QE'
    )
    
    lead_plan = pd.DataFrame(0, index=lead_quarters_index, columns=acquired_customers_plan.columns)
    
    # הלוגיקה הקיימת תעבוד עכשיו על ציר הזמן המורחב
    for q_date, row in acquired_customers_plan.iterrows():
        for c_type in acquired_customers_plan.columns:
            new_cust_count = row[c_type]
            if new_cust_count > 0:
                success_rate = success_rates[c_type] / 100.0
                time_ahead_q = time_aheads_in_quarters[c_type]
                leads_to_contact = np.ceil(new_cust_count / success_rate if success_rate > 0 else 0)
                
                target_period = q_date.to_period('Q') - time_ahead_q
                
                # מציאת הרבעון המתאים בטבלת הלידים המורחבת שלנו
                idx_matches = lead_plan.index[lead_plan.index.to_period('Q') == target_period]
                if len(idx_matches) > 0:
                    lead_plan.loc[idx_matches[0], c_type] += int(leads_to_contact)
                    
    return lead_plan.astype(int)
    # --- END OF CHANGE ---

# --- UI and Logic ---
# !!! שינוי כאן: כותרת גנרית !!!
st.title("Dynamic Multi-Product Business Plan Dashboard")

with st.sidebar:
    st.title("Business Plan Controls")
    with st.expander("Global Parameters (Applied to all products)"):
        st.markdown("**Model Start Date**")
        col1, col2 = st.columns(2)
        with col1:
            model_start_year = st.selectbox("Start Year", options=[2025, 2026, 2027], index=0, key="start_year")
        with col2:
            # הגדרת הרבעון הנוכחי כברירת מחדל
            from datetime import date
            current_quarter = (date.today().month - 1) // 3 + 1
            start_quarter_index = current_quarter - 1 if model_start_year == 2025 else 0
            
            model_start_quarter = st.selectbox("Start Quarter", options=[1, 2, 3, 4], index=start_quarter_index, key="start_quarter")
    
        st.markdown("---") # קו מפריד

    # --- Expander for User & Scenarios ---
    with st.expander("User & Scenarios", expanded=True):
        user_id = st.text_input("Enter your User ID (e.g., email)", key="user_id")
        if user_id and db:
            saved_scenarios = get_user_scenarios(user_id)
            col_load, col_save = st.columns(2)

            # --- עמודה שמאלית: טעינה ומחיקה ---
            # <<< החלף את כל קטע הטעינה והמחיקה שלך בקוד הבא >>>

            with col_load:
                st.subheader("Load or Delete")
                
                # בדיקה אם יש תרחישים שמורים בכלל
                if saved_scenarios:
                    selected_scenario = st.selectbox(
                        "Select scenario",
                        options=saved_scenarios, 
                        index=0, 
                        key="load_scenario_select",
                        label_visibility="collapsed"
                    )
            
                    # --- !!! לוגיקת טעינה מתוקנת !!! ---
                    if st.button("Load Scenario") and selected_scenario:
                        loaded_data = load_scenario_data(user_id, selected_scenario)
                        if loaded_data:
                            # 1. הגדרת מפתחות חיוניים שאסור למחוק בזמן הטעינה
                            keys_to_preserve = [
                                'user_id', 'products', 'results', 'load_scenario_select', 
                                'scenario_name', 'confirm_delete_checkbox', 'new_product_name_input'
                            ]
            
                            # 2. ניקוי כל המפתחות הישנים (ערכי הווידג'טים) מה-session state
                            for key in list(st.session_state.keys()):
                                # בדיקה אם המפתח הוא חיוני או שייך לווידג'ט פנימי של Streamlit
                                is_essential = any(key.startswith(p_key) for p_key in keys_to_preserve)
                                if not is_essential and not key.startswith(('FormSubmitter', '_')):
                                    del st.session_state[key]
                            
                            # 3. איפוס תוצאות וטעינת הנתונים החדשים ל-session state הנקי
                            st.session_state.results = {}
                            for key, value in loaded_data.items():
                                if key not in ['user_id', 'results']: # לא טוענים מחדש את המשתמש או התוצאות
                                    try:
                                        st.session_state[key] = deserialize_from_firestore(value)
                                    except Exception as e:
                                        st.sidebar.error(f"Failed to load key: '{key}'. Error: {e}")
                            
                            st.sidebar.success(f"Scenario '{selected_scenario}' loaded!")
                            st.rerun() # 4. הפעלה מחדש מיידית של האפליקציה
            
                    st.markdown("---")
                    
                    # --- לוגיקת מחיקה (ללא שינוי מהותי) ---
                    if selected_scenario:
                        confirm_delete = st.checkbox(f"Confirm deletion of '{selected_scenario}'", key="confirm_delete_checkbox")
                        if st.button("Delete Scenario", type="primary"):
                            if confirm_delete:
                                if delete_scenario(user_id, selected_scenario):
                                    st.session_state.results = {} # איפוס תוצאות
                                    # ניקוי כל המפתחות כדי לחזור למצב התחלתי נקי
                                    for key in list(st.session_state.keys()):
                                        if key not in ['user_id']: # השאר רק את היוזר
                                            del st.session_state[key]
                                    st.rerun()
                            else:
                                st.warning("Please check the box to confirm.")
                else:
                    st.caption("No scenarios found to load or delete.")
            
            # --- עמודה ימנית: שמירה ---
            with col_save:
                st.subheader("Save New")
                scenario_name_to_save = st.text_input("Save as scenario name:", key="scenario_name")
                if st.button("Save Current") and scenario_name_to_save:
                    if scenario_name_to_save in saved_scenarios:
                        st.error(f"Scenario '{scenario_name_to_save}' already exists.")
                    else:
                        all_inputs = { 'user_id': st.session_state.get('user_id', ''), 'products': st.session_state.get('products', []) }
                        keys_to_exclude = ['results', 'user_id', 'products', 'load_scenario_select', 'scenario_name', 'new_product_name_input', 'confirm_delete_checkbox']
                        for key, value in st.session_state.items():
                            is_excluded = key in keys_to_exclude or key.startswith(('FormSubmitter', '_'))
                            if isinstance(key, str) and not is_excluded:
                                all_inputs[key] = value
                        save_scenario(user_id, scenario_name_to_save, all_inputs)
                        st.rerun()
    # --- התחלה: בלוק AI Analyst (מבוסס Vertex AI) ---
    with st.expander("🤖 AI Analyst (Beta)", expanded=True):
        
        # 1. אתחול ה-API (בודקים סודות Firebase, כי Vertex משתמש בהם)
        if "firebase" not in st.secrets:
            st.error("לא הוגדרו סודות Firebase (נדרש לאימות Vertex AI).")
        else:
            try:
                # אימות ל-Vertex AI משתמש באותם סודות של Firebase
                creds_json = dict(st.secrets.firebase)
                project_id = creds_json.get("project_id")
                
                if not project_id:
                    st.error("project_id חסר בסודות ה-Firebase.")
                else:
                    # --- התחלה: התיקון ---
                    # 1. ניצור אובייקט credentials מלא מתוך הסודות
                    credentials = service_account.Credentials.from_service_account_info(creds_json)
                    
                    # 2. נאתחל את Vertex AI ונעביר לו את ה-credentials במפורש
                    vertexai.init(project=project_id, location="us-central1", credentials=credentials)
    
                    # אתחול המודל עם הגדרת הכלים
                    model = GenerativeModel(
                        "gemini-pro", # זה יעבוד כאן
                        tools=[tools_vertex]
                    )
    
                    # אתחול היסטוריית הצ'אט
                    if "chat_session" not in st.session_state:
                        st.session_state.chat_session = model.start_chat(history=[])
                    
                    # הצגת הודעות קודמות
                    for message in st.session_state.chat_session.history:
                        role = "assistant" if message.role == "model" else message.role
                        with st.chat_message(role):
                            st.markdown(message.parts[0].text)
    
                    # 2. הכנת "הקונטקסט" (זהה לקודם)
                    data_context = "--- נתונים עדכניים ---\n"
                    if "results" in st.session_state and st.session_state.results:
                        data_context += "המשתמש הריץ ניתוח. להלן סיכום התוצאות:\n"
                        try:
                            product_list = [p for p in st.session_state.get('products', []) if p]
                            all_revenues = {p: st.session_state.results[p]['annual_revenue'] for p in product_list if p in st.session_state.results}
                            summary_revenue_df = pd.DataFrame(all_revenues)
                            if pd.api.types.is_datetime64_any_dtype(summary_revenue_df.index):
                                summary_revenue_df.index = summary_revenue_df.index.year
                            data_context += "טבלת סיכום הכנסות שנתיות (כלל המוצרים):\n"
                            data_context += summary_revenue_df.to_markdown() + "\n\n"
                        except Exception as e:
                            pass
                    else:
                        data_context += "המשתמש עדיין לא הריץ ניתוח. הוא נמצא בשלב הגדרת הפרמטרים.\n"
                    data_context += "--- סוף נתונים ---\n"
                    
                    # 3. קבלת שאלה מהמשתמש (זהה)
                    if user_question := st.chat_input("שנה מחיר מוצר 1 ל-20..."):
                        with st.chat_message("user"):
                            st.markdown(user_question)
                        
                        # 4. בניית ההנחיה (Prompt) (זהה)
                        all_setting_keys = [k for k in st.session_state.keys() if isinstance(k, str) and not k.startswith(('_', 'chat_session', 'results', 'messages', 'FormSubmitter'))]
                        
                        prompt_context = f"""
                        אתה עוזר AI שמנהל דשבורד תוכנית עסקית ב-Streamlit.
                        
                        המשימות שלך:
                        1.  **לענות על שאלות:** ענה על שאלות המשתמש לגבי התוצאות (אם קיימות).
                        2.  **לשנות הגדרות:** אם המשתמש מבקש לשנות הגדרה (למשל "שנה מחיר", "הוסף שנה"), עליך להשתמש בכלי `update_setting`.
    
                        מידע חשוב:
                        -   הפורמט של מפתחות הגדרה עבור מוצרים הוא: `key_שםהמוצר`. 
                            לדוגמה, המחיר ההתחלחי של "Product 1" הוא המפתח `ip_unit_Product 1`.
                        -   פרמטרים גלובליים הם פשוטים, למשל `start_year`.
                        
                        ---
                        רשימת מפתחות ההגדרה הקיימים כרגע (לשימושך ב-`setting_key`):
                        {all_setting_keys}
                        ---
                        
                        הנתונים הנוכחיים מהדשבורד:
                        {data_context}
                        ---
                        
                        המשך את השיחה וענה לבקשת המשתמש:
                        """
    
                        # 5. שליחת הבקשה וקבלת תשובה (קצת שונה)
                        try:
                            response = st.session_state.chat_session.send_message(prompt_context + user_question)
                            
                            # 6. בדיקה אם ה-AI רוצה להשתמש בכלי
                            if response.parts[0].function_call:
                                function_call = response.parts[0].function_call
                                function_name = function_call.name
                                
                                if function_name in available_tools:
                                    function_to_call = available_tools[function_name]
                                    function_args = {k: v for k, v in function_call.args.items()} # המרה קטנה
                                    
                                    # --- הפעלת הפונקציה ---
                                    with st.spinner(f"מבצע: {function_name}({function_args.get('setting_key')})..."):
                                        function_response = function_to_call(**function_args)
                                    
                                    # --- שליחת התוצאה חזרה ל-AI ---
                                    # (הפורמט של Vertex קצת שונה)
                                    from vertexai.generative_models import Part
                                    response = st.session_state.chat_session.send_message(
                                        Part.from_function_response(name=function_name, response={"content": function_response})
                                    )
                                    
                                    # הצגת התשובה הסופית של ה-AI
                                    with st.chat_message("assistant"):
                                        st.markdown(response.parts[0].text)
                                    
                                    # --- רענון האפליקציה ---
                                    st.rerun()
    
                                else:
                                    with st.chat_message("assistant"):
                                        st.error(f"ה-AI ניסה לקרוא לפונקציה לא קיימת: {function_name}")
    
                            else:
                                # 7. אם זו תשובה רגילה (טקסט)
                                with st.chat_message("assistant"):
                                    st.markdown(response.parts[0].text)
    
                        except Exception as e:
                            with st.chat_message("assistant"):
                                st.error(f"אירעה שגיאה ב-Vertex AI: {e}")
    
            except Exception as e:
                st.error(f"שגיאה באתחול מודל ה-AI: {e}")
# --- סוף: בלוק AI Analyst (מבוסס Vertex AI) ---                    

    # --- Expander for Managing Products ---
    with st.expander("Manage Products"):
    
        # --- פונקציית עזר פנימית לניקוי ה-state ---
        def remove_product_from_state(product_to_remove):
            # 1. הסרה מרשימת המוצרים הראשית
            if product_to_remove in st.session_state.products:
                st.session_state.products.remove(product_to_remove)
            
            # 2. הסרה מתוצאות הריצה (אם קיימות)
            if product_to_remove in st.session_state.results:
                del st.session_state.results[product_to_remove]
                
            # 3. ניקוי כל המפתחות המשויכים למוצר מה-session_state
            keys_to_delete = []
            for key in st.session_state.keys():
                if isinstance(key, str) and key.endswith(f"_{product_to_remove}"):
                    keys_to_delete.append(key)
            
            for key in keys_to_delete:
                try:
                    del st.session_state[key]
                except KeyError:
                    pass # כבר נמחק, הכל בסדר
        # --- סוף פונקציית העזר ---
    
        # לולאה על עותק של הרשימה כדי לאפשר שינויים תוך כדי
        current_products = st.session_state.get('products', []).copy()
        
        # אם אין מוצרים, הצג הודעה
        if not current_products:
            st.caption("אין מוצרים. הוסף מוצר חדש למטה.")
    
        # הצגת המוצרים הקיימים עם אפשרות מחיקה
        for i, product_name in enumerate(current_products):
            col1, col2 = st.columns([0.8, 0.2]) # עמודה רחבה לשם, צרה לכפתור
            
            with col1:
                # עדכון שם המוצר
                new_name = st.text_input(
                    f"Product {i+1} Name", 
                    value=product_name, 
                    key=f"pname_{i}"
                )
                # אם השם שונה, עדכן אותו ברשימה הראשית
                if new_name != product_name and new_name:
                    st.session_state.products[i] = new_name
                    # (כאן אפשר להוסיף לוגיקה מורכבת יותר של שינוי שם גם ב-keys)
                    # כרגע נשאיר פשוט - המשתמש יצטרך להריץ מחדש
            
            with col2:
                st.write("") # טריק קטן ליישור אנכי
                st.write("")
                st.button(
                    "✖️", # כפתור הסרה
                    key=f"remove_prod_{i}", 
                    on_click=remove_product_from_state, # קריאה לפונקציית הניקוי
                    args=(product_name,), # הפרמטר שיועבר לפונקציה
                    help=f"הסר את {product_name}"
                )
    
        st.markdown("---")
        
        # לוגיקה להוספת מוצר חדש (זהה לקודם)
        new_product_name = st.text_input("New Product Name", key="new_product_name_input")
        if st.button("Add Product") and new_product_name:
            if new_product_name not in st.session_state.products:
                st.session_state.products.append(new_product_name)
                st.rerun() # רענון מיידי כדי שהמוצר החדש יופיע
            else:
                st.warning("Product name already exists.")

    # --- Expander for Lead Generation Parameters ---
    with st.expander("Lead Generation Parameters (Global)"):
        lead_params = { 'success_rates': {}, 'time_aheads_in_quarters': {} }
        customer_types_for_leads = ['Medium', 'Large', 'Global']
        sr_defaults = {'Medium': 50, 'Large': 40, 'Global': 30}
        ta_defaults = {'Medium': 3, 'Large': 4, 'Global': 6}
        for c_type in customer_types_for_leads:
            sr_key = f'sr_{c_type}'
            ta_key = f'ta_{c_type}'
            lead_params['success_rates'][c_type] = st.slider(f'Success Rate (%) - {c_type}', 0, 100, st.session_state.get(sr_key, sr_defaults[c_type]), key=sr_key)
            lead_params['time_aheads_in_quarters'][c_type] = st.slider(f'Time Ahead (Quarters) - {c_type}', 1, 12, st.session_state.get(ta_key, ta_defaults[c_type]), key=ta_key)
        
    # --- Product-specific Inputs ---
    product_inputs = {}
    for product in st.session_state.get('products', []).copy():
        st.header(product)
        product_inputs[product] = {}
        
        with st.expander(f"1. Initial Customer Value & Launch Year", expanded=False):
            
            product_inputs[product]['launch_year'] = st.selectbox(
                "Launch Year", 
                options=[2025, 2026, 2027, 2028, 2029, 2030], 
                index=0, 
                key=f'launch_year_{product}'
            )
            st.markdown("---")
            
            # !!! שינוי כאן: טונות -> יחידות !!!
            product_inputs[product]['is_m'] = st.number_input('Initial Units/Customer - Medium:', 0.0, value=st.session_state.get(f'is_m_{product}', 1.5), step=0.1, key=f'is_m_{product}')
            product_inputs[product]['is_l'] = st.number_input('Initial Units/Customer - Large:', 0.0, value=st.session_state.get(f'is_l_{product}', 10.0), step=1.0, key=f'is_l_{product}')
            product_inputs[product]['is_g'] = st.number_input('Initial Units/Customer - Global:', 0.0, value=st.session_state.get(f'is_g_{product}', 40.0), step=2.0, key=f'is_g_{product}')
            
        with st.expander(f"2. Customer Value Growth", expanded=False):
            product_inputs[product]['market_gr'] = st.slider('Annual Market Growth Rate (%):', 0.0, 20.0, st.session_state.get(f'mgr_{product}', 6.4), 0.1, key=f'mgr_{product}')
            product_inputs[product]['pen_y1'] = st.slider('Penetration Rate Year 1 (%):', 1.0, 20.0, st.session_state.get(f'pen_y1_{product}', 7.5), 0.1, key=f'pen_y1_{product}')
            # !!! שינוי כאן: טונות -> יחידות !!!
            product_inputs[product]['tt_m'] = st.number_input('Target Units/Cust Year 5 - Medium:', 0.0, value=st.session_state.get(f'tt_m_{product}', 89.0), key=f'tt_m_{product}')
            product_inputs[product]['tt_l'] = st.number_input('Target Units/Cust Year 5 - Large:', 0.0, value=st.session_state.get(f'tt_l_{product}', 223.0), key=f'tt_l_{product}')
            product_inputs[product]['tt_g'] = st.number_input('Target Units/Cust Year 5 - Global:', 0.0, value=st.session_state.get(f'tt_g_{product}', 536.0), key=f'tt_g_{product}')
            
        with st.expander(f"3. Revenue Targets & Sales Strategy", expanded=False):
            st.markdown("**Target Annual Revenue ($)**")
            default_revenues = [250000, 2700000, 5500000, 12000000, 32000000, 40000000]
            rev_targets = []
            for i in range(6):
                year_num = i + 1
                rev_val = st.number_input(
                    f'Year {year_num}:', 
                    min_value=0, 
                    max_value=100_000_000, 
                    value=st.session_state.get(f'rev_y{year_num}_{product}', default_revenues[i]), 
                    step=10000, 
                    key=f'rev_y{year_num}_{product}'
                )
                rev_targets.append(rev_val)
            product_inputs[product]['annual_rev_targets'] = rev_targets
            st.markdown("---")
            st.markdown("**Sales Focus (%)**")
            product_inputs[product]['f_m'] = st.slider('Medium:', 0, 100, st.session_state.get(f'f_m_{product}', 50), 5, key=f'f_m_{product}')
            product_inputs[product]['f_l'] = st.slider('Large:', 0, 100, st.session_state.get(f'f_l_{product}', 30), 5, key=f'f_l_{product}')
            product_inputs[product]['f_g'] = st.slider('Global:', 0, 100, st.session_state.get(f'f_g_{product}', 20), 5, key=f'f_g_{product}')
            
        with st.expander(f"4. Pricing Assumptions", expanded=False):
            # !!! שינוי כאן: ק"ג -> יחידה + שינוי מפתח !!!
            product_inputs[product]['ip_unit'] = st.number_input('Initial Price per Unit ($):', 0.0, value=st.session_state.get(f'ip_unit_{product}', 18.0), step=0.5, key=f'ip_unit_{product}')
            product_inputs[product]['pdr'] = st.slider('Quarterly Price Decay (%):', 0.0, 10.0, st.session_state.get(f'pdr_{product}', 3.65), 0.05, key=f'pdr_{product}')
            product_inputs[product]['price_floor'] = st.number_input('Minimum Price ($):', 0.0, value=st.session_state.get(f'price_floor_{product}', 14.0), step=0.5, key=f'price_floor_{product}')
            
        with st.expander(f"5. Production Costs ($/Unit)", expanded=False):
            # !!! שינוי כאן: טונות -> יחידות !!!
            st.markdown("Define cost based on quarterly production volume (in Units)")
            
            # הגדרת ערכי ברירת המחדל מהתמונה
            default_quantities = [10, 20, 40, 100, 200, 1500]
            default_costs = [15.32, 13.14, 10.73, 8.46, 8.37, 7.43]
            
            cost_quantities = []
            cost_values = []
            
            cols = st.columns(6)
            for i in range(6):
                with cols[i]:
                    # !!! שינוי כאן: מפתחות וטקסט עבור יחידות !!!
                    q_key = f'cost_u_{i}_{product}' # u for units
                    c_key = f'cost_c_{i}_{product}'
                    
                    qty = st.number_input(f"Units {i+1}", value=st.session_state.get(q_key, default_quantities[i]), key=q_key)
                    cost = st.number_input(f"Cost {i+1}", value=st.session_state.get(c_key, default_costs[i]), format="%.2f", key=c_key)
                    cost_quantities.append(qty)
                    cost_values.append(cost)
            
            # !!! שינוי כאן: טונות -> יחידות !!!
            product_inputs[product]['cost_quantities_u'] = cost_quantities
            product_inputs[product]['cost_values_per_unit'] = cost_values
        
# --- Run Button ---
run_button = st.sidebar.button("Run Full Analysis", use_container_width=True)

# --- App Logic and Display ---
if run_button:
    results_data = {}
    # Use a copy of the list to avoid issues if it's modified
    for product in st.session_state.get('products', []).copy():
        res = calculate_plan(**product_inputs[product], global_start_year=model_start_year, global_start_quarter=model_start_quarter)
        if res.get("error"):
            st.error(f"Error for {product}: {res['error']}"); st.stop()
        
        final_cumulative = res["cumulative_customers"].round().astype(int)
        acquired_customers = final_cumulative.diff(axis=0).fillna(final_cumulative.iloc[0]).clip(lower=0).astype(int)
        
        res['acquired_customers_plan'] = acquired_customers
        res['cumulative_customers'] = final_cumulative
        res['lead_plan'] = create_lead_plan(acquired_customers, **lead_params)
        results_data[product] = res
    st.session_state.results = results_data

# <<< החלף את כל קוד התצוגה שלך (מ-if st.session_state.results: ועד הסוף) בקוד המלא הבא >>>

if st.session_state.results:
    results = st.session_state.results
    # מסנן החוצה מוצרים עם שם ריק כדי למנוע שגיאות
    product_list = [p for p in st.session_state.get('products', []) if p]
    tabs = st.tabs([*product_list, "Overall Summary"])
    
    # --- הגדרת תאריכי תצוגה לפי הבקשה ---
    lead_display_start_date = pd.Timestamp('2025-01-01')
    main_display_start_date = pd.Timestamp('2025-07-01') # התחלת תצוגה מרבעון 3, 2025

    # --- פונקציית עזר לשימוש חוזר (מוגדרת פעם אחת מחוץ ללולאה) ---
    def format_quarterly_cols(df_transposed):
        df_copy = df_transposed.copy()
        if not df_copy.empty:
            df_copy.columns = [f"{c.year}-Q{c.quarter}" for c in df_copy.columns]
        return df_copy
    
    # --- לולאה להצגת התוצאות בכל לשונית של מוצר ---
    for i, product_name in enumerate(product_list):
        with tabs[i]:
            st.header(f"Results for {product_name}")
            
            # טבלת סיכום רווחיות שנתית
            st.subheader("Profitability Summary (Yearly)")
            profit_summary_df = pd.DataFrame({
                "Total Revenue": results[product_name]['annual_revenue'],
                "Total Production Cost": results[product_name]['total_production_cost_q'].resample('YE').sum(),
                "Total Profit": results[product_name]['profit_q'].resample('YE').sum()
            })
            profit_summary_df["Profit Margin (%)"] = (profit_summary_df["Total Profit"] / profit_summary_df["Total Revenue"].replace(0, np.nan)).fillna(0) * 100
            if pd.api.types.is_datetime64_any_dtype(profit_summary_df.index):
                profit_summary_df.index = profit_summary_df.index.year
            st.dataframe(profit_summary_df.style.format({
                "Total Revenue": "${:,.0f}",
                "Total Production Cost": "${:,.0f}",
                "Total Profit": "${:,.0f}",
                "Profit Margin (%)": "{:.1f}%"
            }))
            st.markdown("---")
            
            # --- סינון נתונים לתצוגה לפי התאריכים החדשים ---
            leads_to_display = results[product_name]["lead_plan"][results[product_name]["lead_plan"].index >= lead_display_start_date]
            acquired_to_display = results[product_name]["acquired_customers_plan"][results[product_name]["acquired_customers_plan"].index >= main_display_start_date]
            cumulative_to_display = results[product_name]["cumulative_customers"][results[product_name]["cumulative_customers"].index >= main_display_start_date]
            
            revenue_by_type_to_display = results[product_name]["revenue_by_type_q"][results[product_name]["revenue_by_type_q"].index >= main_display_start_date]
            # !!! שינוי כאן: טונות -> יחידות !!!
            units_by_type_to_display = results[product_name]["units_by_type_q"][results[product_name]["units_by_type_q"].index >= main_display_start_date]
        
            # --- הצגת שאר התוצאות ---
            st.subheader("Lead Generation")
            st.markdown("#### Table 0: Recommended Lead Contact Plan")
            lead_plan_display = format_quarterly_cols(leads_to_display.T)
            st.dataframe(lead_plan_display.style.format("{:d}"))
            
            st.markdown("##### Chart 0: Yearly Lead Contact Plan")
            leads_for_chart0 = leads_to_display[leads_to_display.index.year != 2030]
            fig0 = create_yearly_bar_chart(df_quarterly=leads_for_chart0, title=f"Leads to Contact per Year - {product_name}", y_axis_label="Number of Leads to Contact")
            st.pyplot(fig0)
            st.markdown("---")
            
            st.subheader("Action Plan & Outcomes")
        
            st.markdown("#### Table A: Revenue per Customer Type ($)")
            revenue_by_type_display = format_quarterly_cols(revenue_by_type_to_display.T)
            st.dataframe(revenue_by_type_display.style.format("${:,.0f}"))
        
            # !!! שינוי כאן: טונות -> יחידות !!!
            st.markdown("#### Table B: Units Sold per Customer Type")
            units_by_type_display = format_quarterly_cols(units_by_type_to_display.T)
            st.dataframe(units_by_type_display.style.format("{:,.2f}"))
            st.markdown("---")
            
            st.markdown("#### Table 1: Acquired New Customers per Quarter")
            acquired_customers_display = format_quarterly_cols(acquired_to_display.T)
            st.dataframe(acquired_customers_display.style.format("{:d}"))
            
            st.markdown("##### Chart 1: Yearly Acquired New Customers")
            fig1 = create_yearly_bar_chart(df_quarterly=acquired_to_display, title=f"Acquired New Customers per Year - {product_name}", y_axis_label="Number of New Customers")
            st.pyplot(fig1)
            st.markdown("---")
            
            st.markdown("#### Table 2: Cumulative Number of Customers (Quarterly)")
            cum_cust_display = format_quarterly_cols(cumulative_to_display.T)
            st.dataframe(cum_cust_display.style.format("{:,d}"))
            
            st.markdown("##### Chart 2: Cumulative Customers (End of Year)")
            fig2 = create_yearly_bar_chart(df_quarterly=cumulative_to_display, title=f"Cumulative Customers at Year End - {product_name}", y_axis_label="Total Number of Customers", is_cumulative=True)
            st.pyplot(fig2)
        
            st.markdown("---")
            st.subheader("Interactive Analysis: Customer & YTD Revenue Mix")
            quarter_options = results[product_name]['cumulative_customers'].index
            selected_quarter = st.selectbox(
                "Select a Quarter to Analyze",
                options=quarter_options,
                format_func=lambda d: f"{d.year}-Q{d.quarter}",
                key=f'pie_select_{product_name}'
            )
            if selected_quarter:
                cust_data_for_quarter = results[product_name]['cumulative_customers'].loc[selected_quarter]
                selected_year = selected_quarter.year
                start_of_year = pd.Timestamp(f'{selected_year}-01-01')
                
                all_rev_data = results[product_name]['revenue_by_type_q'] 
                
                ytd_rev_df = all_rev_data[(all_rev_data.index >= start_of_year) & (all_rev_data.index <= selected_quarter)]
                ytd_rev_per_segment = ytd_rev_df.sum()
                non_zero_cust_data = cust_data_for_quarter[cust_data_for_quarter > 0]
                if not non_zero_cust_data.empty:
                    fig_pie, ax_pie = plt.subplots(figsize=(9, 6))
                    total_customers = non_zero_cust_data.sum()
                    pie_labels = []
                    for segment, count in non_zero_cust_data.items():
                        percentage = (count / total_customers) * 100
                        revenue = ytd_rev_per_segment.get(segment, 0)
                        rev_text = f"${revenue/1_000_000:.2f}M" if revenue >= 1_000_000 else f"${revenue/1_000:,.0f}K"
                        pie_labels.append(f"{segment}\n{percentage:.1f}%\n(YTD: {rev_text})")
                    colors = sns.color_palette('crest', n_colors=len(non_zero_cust_data))
                    wedges, texts = ax_pie.pie(non_zero_cust_data, labels=pie_labels, colors=colors, startangle=90, wedgeprops=dict(width=0.4, edgecolor='w'), textprops={'fontsize': 11})
                    ax_pie.set_title(f"Customer Mix & YTD Revenue Contribution for {selected_quarter.year}-Q{selected_quarter.quarter}", fontsize=16, weight='bold')
                    st.pyplot(fig_pie)
                else:
                    st.info(f"No cumulative customers found for {selected_quarter.year}-Q{selected_quarter.quarter}.")
            
            st.markdown("---")
            
            actual_revenue_series = results[product_name]['annual_revenue'].copy()
            if pd.api.types.is_datetime64_any_dtype(actual_revenue_series.index):
                actual_revenue_series.index = actual_revenue_series.index.year
            
            validation_df = pd.DataFrame({
                'Target Revenue': results[product_name]['annual_revenue_targets'], 
                'Actual Revenue': actual_revenue_series
            })
            validation_df.index.name = "Year"
            results[product_name]['validation_df'] = validation_df
            st.markdown("#### Table 3: Target vs. Actual Revenue")
            st.dataframe(validation_df.style.format({'Target Revenue': "${:,.0f}", 'Actual Revenue': "${:,.0f}"}))
            st.markdown("#### Chart: Sales Income ($)")
            
            plot_data = validation_df[['Actual Revenue']].reset_index()
            
            fig, ax = plt.subplots(figsize=(14, 7))
            
            barplot = sns.barplot(data=plot_data, x='Year', y='Actual Revenue', ax=ax, color='teal')
            
            ax.set_title(f'Sales Income ($) - {product_name}', fontsize=18, weight='bold')
            ax.get_yaxis().set_major_formatter(plt.FuncFormatter(lambda x, p: f"${x/1_000_000:.1f}M"))
            ax.set_xlabel("Year", fontsize=12)
            ax.set_ylabel("Sales Income ($)", fontsize=12)
            for container in barplot.containers:
                ax.bar_label(container, fmt='${:,.0f}', padding=3, fontsize=10)
            st.pyplot(fig)
            with st.expander("View Underlying Assumptions"):
                # !!! שינוי כאן: טונות -> יחידות !!!
                units_per_customer_df = results[product_name].get('units_per_customer')
                pen_rate_df = results[product_name].get('pen_rate_df')
                if units_per_customer_df is not None:
                    st.markdown("#### Table 4: Annual Units per Single Customer (Target-Driven)")
                    st.dataframe(units_per_customer_df.T.style.format("{:,.2f}"))
                if pen_rate_df is not None:
                    st.markdown("#### Table 5: Generated Penetration Rates to Meet Target (%)")
                    st.dataframe((pen_rate_df.T*100).style.format("{:,.1f}%"))
            st.markdown("---")
            if product_name: 
                col1, col2 = st.columns(2)
                with col1:
                    excel_product_data = to_excel({product_name: results[product_name]})
                    if excel_product_data:
                        st.download_button(label=f"📥 Download {product_name} to Excel", data=excel_product_data, file_name=f"{product_name}_Report.xlsx", use_container_width=True)
                with col2:
                    ppt_product_data = create_product_presentation(product_name, results[product_name])
                    if ppt_product_data:
                        st.download_button(label=f"📊 Download {product_name} Presentation", data=ppt_product_data, file_name=f"{product_name}_Presentation.pptx", use_container_width=True)

    # --- לשונית הסיכום הכללי מתחילה כאן (אחרי לולאת המוצרים) ---
    # <<< החלף את כל הקוד של לשונית הסיכום שלך בקוד המלא והמשולב הבא >>>
    
    with tabs[-1]:
        st.header("Overall Summary (All Products)")
    
        # --- חלק 1: הצגת כל טבלאות הסיכום (הקוד המקורי שלך, ללא שינוי) ---
        # --- סיכום הכנסות שנתי ---
        summary_revenue_list = [results[p]['annual_revenue'] for p in product_list if p in results]
        if summary_revenue_list:
            summary_revenue_df = pd.concat(summary_revenue_list, axis=1).sum(axis=1).to_frame(name="Total Revenue")
            if pd.api.types.is_datetime64_any_dtype(summary_revenue_df.index):
                summary_revenue_df.index = summary_revenue_df.index.year
            st.markdown("#### Summary: Total Revenue per Year")
            st.dataframe(summary_revenue_df.style.format("${:,.0f}"))
    
        # --- סיכום הכנסות רבעוני ---
        quarterly_revenues_by_product = {p: results[p]['revenue_by_type_q'].sum(axis=1) for p in product_list if p in results}
        if quarterly_revenues_by_product:
            summary_quarterly_rev_df = pd.DataFrame(quarterly_revenues_by_product)
            summary_quarterly_rev_df['Total'] = summary_quarterly_rev_df.sum(axis=1)
            summary_quarterly_rev_to_display = summary_quarterly_rev_df[summary_quarterly_rev_df.index >= main_display_start_date]
            st.markdown("#### Summary: Quarterly Revenue by Product")
            st.dataframe(format_quarterly_cols(summary_quarterly_rev_to_display.T).style.format("${:,.0f}"))
    
        # !!! שינוי כאן: טונות -> יחידות !!!
        # --- סיכום יחידות רבעוני ---
        quarterly_units_by_product = {p: results[p]['units_by_type_q'].sum(axis=1) for p in product_list if p in results}
        if quarterly_units_by_product:
            summary_quarterly_units_df = pd.DataFrame(quarterly_units_by_product)
            summary_quarterly_units_df['Total'] = summary_quarterly_units_df.sum(axis=1)
            summary_quarterly_units_to_display = summary_quarterly_units_df[summary_quarterly_units_df.index >= main_display_start_date]
            st.markdown("#### Summary: Quarterly Units Sold by Product")
            st.dataframe(format_quarterly_cols(summary_quarterly_units_to_display.T).style.format("{:,.2f}"))
    
        # --- סיכום רווח רבעוני ---
        quarterly_profit_by_product = {p: results[p]['profit_q'] for p in product_list if p in results}
        if quarterly_profit_by_product:
            summary_quarterly_profit_df = pd.DataFrame(quarterly_profit_by_product)
            summary_quarterly_profit_df['Total'] = summary_quarterly_profit_df.sum(axis=1)
            summary_quarterly_profit_to_display = summary_quarterly_profit_df[summary_quarterly_profit_df.index >= main_display_start_date]
            st.markdown("#### Summary: Quarterly Profit by Product")
            st.dataframe(format_quarterly_cols(summary_quarterly_profit_to_display.T).style.format("${:,.0f}"))
    
        # --- סיכום לקוחות מצטבר ---
        summary_customers_list = [results[p]['cumulative_customers'] for p in product_list if p in results]
        if summary_customers_list:
            summary_customers_total_q_raw = pd.concat(summary_customers_list, axis=1).sum(axis=1)
            summary_customers_to_display = summary_customers_total_q_raw[summary_customers_total_q_raw.index >= main_display_start_date]
            summary_customers_display_T = summary_customers_to_display.to_frame(name="Total Customers").T
            st.markdown("#### Summary: Total Cumulative Customers (Quarterly)")
            st.dataframe(format_quarterly_cols(summary_customers_display_T).style.format("{:,d}"))
        
    
        # --- !!! חלק 2: הצגת התרשים הדינמי (הקוד החדש) !!! ---
        all_revenues = {p: results[p]['annual_revenue'] for p in product_list if p in results}
        
        if all_revenues:
            summary_plot_df = pd.DataFrame(all_revenues)
            if pd.api.types.is_datetime64_any_dtype(summary_plot_df.index):
                summary_plot_df.index = summary_plot_df.index.year
    
            st.markdown("---")
            st.markdown("#### Chart: Total Revenue Breakdown by Product")
    
            # בורר המצבים החדש לבחירת סוג התרשים
            chart_type = st.radio(
                "Select Chart Style:",
                options=["Interactive (Plotly)", "Stacked Bar", "Grouped Bar"],
                horizontal=True,
                label_visibility="collapsed"
            )
    
            if chart_type == "Stacked Bar":
                fig_to_show = create_stacked_bar_chart(summary_plot_df)
                st.pyplot(fig_to_show)
            
            elif chart_type == "Grouped Bar":
                fig_to_show = create_grouped_bar_chart(summary_plot_df)
                st.pyplot(fig_to_show)
                
            elif chart_type == "Interactive (Plotly)":
                fig_to_show = create_interactive_plotly_chart(summary_plot_df)
                st.plotly_chart(fig_to_show, use_container_width=True)
    
        # --- חלק 3: כפתורי הורדה (הקוד המקורי שלך, ללא שינוי) ---
        st.markdown("---")
        if 'summary_revenue_df' in locals() and 'summary_customers_total_q_raw' in locals():
            col1, col2, col3 = st.columns(3)
            summary_for_excel = {"summary_revenue": summary_revenue_df, "summary_customers_raw": summary_customers_total_q_raw}
            with col1:
                excel_summary_data = to_excel({"summary": summary_for_excel, **results})
                if excel_summary_data:
                    st.download_button(label="📥 Download Summary to Excel", data=excel_summary_data, file_name="Overall_Summary_Report.xlsx", use_container_width=True)
            with col2:
                ppt_summary_data = create_summary_presentation(summary_for_excel, results)
                if ppt_summary_data:
                    st.download_button(label="📊 Download Summary Presentation", data=ppt_summary_data, file_name="Overall_Summary_Presentation.pptx", use_container_width=True)
            with col3:
                pdf_data = to_pdf(results)
                if pdf_data:
                    st.download_button(label="📄 Download Full PDF Report", data=pdf_data, file_name="Full_Analysis_Report.pdf", use_container_width=True)

if not st.session_state.results:
    st.info("Set your parameters in the sidebar and click 'Run Full Analysis' to see the results.")
